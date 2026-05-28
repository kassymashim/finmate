"""Generate PowerPoint presentation from FinMate content."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors matching our HTML theme
BG_DARK = RGBColor(0x0F, 0x11, 0x17)
BG_CARD = RGBColor(0x1E, 0x22, 0x35)
TEXT_WHITE = RGBColor(0xF1, 0xF5, 0xF9)
TEXT_GRAY = RGBColor(0x94, 0xA3, 0xB8)
TEXT_LIGHT = RGBColor(0xCB, 0xD5, 0xE1)
ACCENT_PURPLE = RGBColor(0x81, 0x8C, 0xF8)
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)
ACCENT_INDIGO = RGBColor(0x63, 0x66, 0xF1)


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=TEXT_WHITE, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tf


def add_bullet_points(slide, left, top, width, height, items, font_size=14, color=TEXT_LIGHT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"→  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
    return tf


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # Blank layout

    # =========== SLIDE 1: Title ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text_box(slide, Inches(1), Inches(2.2), Inches(11.3), Inches(1.5),
                 "FinMate", font_size=60, bold=True, color=ACCENT_INDIGO,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.5), Inches(11.3), Inches(0.8),
                 "AI-Powered Personal Finance Assistant", font_size=24,
                 color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(5.0), Inches(11.3), Inches(0.8),
                 "LLM Engineering Course — Final Project  |  May 2026\ngithub.com/kassymashim/finmate",
                 font_size=14, color=RGBColor(0x64, 0x74, 0x8B), alignment=PP_ALIGN.CENTER)

    # =========== SLIDE 2: Problem ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                 "The Problem", font_size=36, bold=True, color=TEXT_WHITE)
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.7),
                 '"65% of people don\'t know how much they spent last month."',
                 font_size=18, color=ACCENT_PURPLE)
    add_text_box(slide, Inches(0.8), Inches(2.2), Inches(5), Inches(0.5),
                 "Pain Points", font_size=18, bold=True, color=ACCENT_PURPLE)
    add_bullet_points(slide, Inches(0.8), Inches(2.7), Inches(5.5), Inches(3.5), [
        "Manual expense tracking is tedious → people give up",
        "Banking apps show data but don't explain patterns",
        "Generic budgeting advice ignores individual context",
        "Financial advisors are expensive ($150-300/hour)"
    ])
    add_text_box(slide, Inches(7), Inches(2.2), Inches(5), Inches(0.5),
                 "Current Alternatives", font_size=18, bold=True, color=ACCENT_PURPLE)
    add_bullet_points(slide, Inches(7), Inches(2.7), Inches(5.5), Inches(3.5), [
        "Spreadsheets — boring, requires discipline",
        "Banking apps — raw data, no guidance",
        "Budgeting apps — rules without personalization",
        "Financial advisors — costly, inaccessible"
    ])

    # =========== SLIDE 3: Solution ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                 "The Solution — FinMate", font_size=36, bold=True, color=TEXT_WHITE)
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.7),
                 "An AI assistant that turns transaction data into personalized, actionable financial advice.",
                 font_size=18, color=TEXT_LIGHT)
    features = [
        ("📊 Smart Dashboard", "Visual spending analytics with category breakdowns, trends, and monthly comparisons"),
        ("💬 AI Advisor", "Conversational chat that knows YOUR spending patterns and gives personalized advice"),
        ("⚡ Quick Logging", 'Type "Chipotle $15" → instant categorization and dashboard update'),
        ("📷 Receipt OCR", "Photo receipts analyzed by GPT-4o Vision — auto-extracted and categorized"),
        ("🎯 Goal Planning", "Milestone-based savings plans with timelines and action items"),
        ("🛡️ Privacy First", "PII auto-masking, prompt injection detection, domain guardrails"),
    ]
    for i, (title, desc) in enumerate(features):
        col = i % 3
        row = i // 3
        left = Inches(0.8 + col * 4.1)
        top = Inches(2.3 + row * 2.2)
        add_text_box(slide, left, top, Inches(3.8), Inches(0.5),
                     title, font_size=14, bold=True, color=ACCENT_PURPLE)
        add_text_box(slide, left, top + Inches(0.45), Inches(3.8), Inches(1.2),
                     desc, font_size=12, color=TEXT_GRAY)

    # =========== SLIDE 4: Demo ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text_box(slide, Inches(1), Inches(0.5), Inches(11.3), Inches(0.8),
                 "Live Demo", font_size=36, bold=True, color=TEXT_WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(1.5), Inches(5), Inches(0.5),
                 "Demo Flow", font_size=18, bold=True, color=ACCENT_PURPLE)
    add_bullet_points(slide, Inches(1.5), Inches(2.0), Inches(5), Inches(4), [
        "Dashboard — spending breakdown",
        "Month filter — compare periods",
        'Quick add — "Starbucks $6"',
        "AI Chat — spending analysis",
        "AI Chat — savings plan",
        "LangSmith — real traces"
    ])
    add_text_box(slide, Inches(7), Inches(1.5), Inches(5), Inches(0.5),
                 "Key Questions to Answer", font_size=18, bold=True, color=ACCENT_PURPLE)
    add_bullet_points(slide, Inches(7), Inches(2.0), Inches(5), Inches(4), [
        '"Where is my money going?"',
        '"How can I save $15,000?"',
        '"Am I overspending on food?"',
        '"Create a budget for me"'
    ])

    # =========== SLIDE 5: Architecture ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                 "Architecture Overview", font_size=36, bold=True, color=TEXT_WHITE)
    arch_text = (
        "User → Next.js Frontend → FastAPI Backend\n"
        "                                    ↓\n"
        "                    Guardrails (PII filter, injection detection)\n"
        "                                    ↓\n"
        "                    LangGraph Orchestration\n"
        "                    ├── Classify Intent (GPT-4o, temp=0)\n"
        "                    ├── Retrieve Context (ChromaDB RAG, MMR)\n"
        "                    └── Route to Specialist:\n"
        "                        ├── Budget Analysis\n"
        "                        ├── Savings Advice\n"
        "                        ├── Spending Insights\n"
        "                        ├── Goal Setting (human-in-the-loop)\n"
        "                        └── General Q&A\n"
        "                                    ↓\n"
        "                    LangSmith (tracing all calls)"
    )
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5),
                 arch_text, font_size=13, color=TEXT_GRAY)

    # =========== SLIDE 6: LangGraph ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                 "LangGraph Workflow", font_size=36, bold=True, color=TEXT_WHITE)
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5), Inches(0.5),
                 "Why LangGraph?", font_size=18, bold=True, color=ACCENT_PURPLE)
    add_bullet_points(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4), [
        "Granular control over routing logic",
        "Native conditional edges + state management",
        "Single-agent, multiple-roles pattern",
        "Explicit graph = easier debugging",
        "Considered CrewAI — too much overhead for our use case"
    ])
    add_text_box(slide, Inches(7), Inches(1.5), Inches(5), Inches(0.5),
                 "Multi-step Flow", font_size=18, bold=True, color=ACCENT_PURPLE)
    add_bullet_points(slide, Inches(7), Inches(2.0), Inches(5.5), Inches(4), [
        "Step 1: Intent Classification (5 categories)",
        "Step 2: RAG Retrieval (enhanced query)",
        "Step 3: Conditional Routing (5 specialist nodes)",
        "Step 4: Response Generation (context-aware)",
        "Step 5: Human-in-the-loop (goal confirmation)"
    ])

    # =========== SLIDE 7: RAG + MCP ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                 "RAG Pipeline & MCP Server", font_size=36, bold=True, color=TEXT_WHITE)
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5), Inches(0.5),
                 "RAG Pipeline", font_size=18, bold=True, color=ACCENT_PURPLE)
    add_bullet_points(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5), [
        "Knowledge: 4 curated finance documents",
        "Chunking: 800 chars, 200 overlap, MD-aware",
        "Embeddings: text-embedding-3-small",
        "Vector DB: ChromaDB (local, persistent)",
        "Retrieval: MMR (λ=0.7) for diversity",
        "Why ChromaDB: Zero cost, sufficient scale"
    ])
    add_text_box(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.5),
                 "Custom MCP Server (3 Tools)", font_size=18, bold=True, color=ACCENT_PURPLE)
    add_bullet_points(slide, Inches(7), Inches(2.0), Inches(5.5), Inches(4.5), [
        "categorize_expenses — Auto-categorize by merchant keywords",
        "calculate_budget — 50/30/20 allocation with adjustments",
        "get_savings_plan — Milestone plan with timeline & actions"
    ])

    # =========== SLIDE 8: Multimodality ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                 "Multimodality & Guardrails", font_size=36, bold=True, color=TEXT_WHITE)
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5), Inches(0.5),
                 "GPT-4o Vision — Receipt OCR", font_size=18, bold=True, color=ACCENT_PURPLE)
    add_bullet_points(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4), [
        "Upload receipt/bank statement photo",
        "Extracts: merchant, date, items, total",
        "Auto-categorizes the expense",
        "Structured JSON output for dashboard",
        "Supports receipts, invoices, statements"
    ])
    add_text_box(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.5),
                 "Guardrails (Bonus)", font_size=18, bold=True, color=ACCENT_PURPLE)
    add_bullet_points(slide, Inches(7), Inches(2.0), Inches(5.5), Inches(4), [
        "PII Detection: Credit cards, SSNs auto-masked",
        "Injection Detection: Blocks prompt overrides",
        "Domain Enforcement: Stays within finance",
        "Output Validation: Adds disclaimers when needed"
    ])

    # =========== SLIDE 9: Evals ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                 "Evaluation Results", font_size=36, bold=True, color=TEXT_WHITE)
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.5),
                 "Golden Dataset: 35 examples across 5 intent categories",
                 font_size=16, color=TEXT_LIGHT)
    # Metrics
    metrics = [("80%", "Intent Accuracy"), ("3.93/5", "Response Relevance"), ("92.2%", "Keyword Faithfulness")]
    for i, (val, label) in enumerate(metrics):
        left = Inches(1.5 + i * 3.8)
        add_text_box(slide, left, Inches(2.2), Inches(3), Inches(0.8),
                     val, font_size=40, bold=True, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left, Inches(3.0), Inches(3), Inches(0.5),
                     label, font_size=14, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)
    # A/B table
    add_text_box(slide, Inches(0.8), Inches(3.9), Inches(11), Inches(0.5),
                 "A/B Test: GPT-4o vs GPT-4o-mini", font_size=18, bold=True, color=ACCENT_PURPLE)
    ab_text = (
        "Intent Accuracy:   GPT-4o: 80%     GPT-4o-mini: ~72%    → Winner: GPT-4o\n"
        "Latency:           GPT-4o: 5.56s   GPT-4o-mini: ~2.8s   → Winner: GPT-4o-mini\n"
        "Cost per query:    GPT-4o: $0.01   GPT-4o-mini: $0.001  → Winner: GPT-4o-mini\n\n"
        "Decision: GPT-4o for responses (quality), GPT-4o-mini for intent classification (speed/cost)"
    )
    add_text_box(slide, Inches(0.8), Inches(4.4), Inches(11.5), Inches(2.5),
                 ab_text, font_size=13, color=TEXT_LIGHT)

    # =========== SLIDE 10: Hyperparameters ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                 "LLM Choice & Hyperparameters", font_size=36, bold=True, color=TEXT_WHITE)
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5), Inches(0.5),
                 "Why GPT-4o?", font_size=18, bold=True, color=ACCENT_PURPLE)
    add_bullet_points(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(2.5), [
        "Best vision capabilities (receipt OCR)",
        "Superior instruction following",
        "Consistent structured JSON output",
        "Acceptable latency for chat (~5.5s)",
        "Fallback: GPT-4o-mini if rate limited"
    ])
    add_text_box(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.5),
                 "Hyperparameter Experiments", font_size=18, bold=True, color=ACCENT_PURPLE)
    hp_text = (
        "Temperature:  Tested 0, 0.3, 0.7, 1.0  →  Selected: 0.3\n"
        "Top-P:        Tested 0.8, 0.9, 1.0      →  Selected: 0.9\n"
        "Max Tokens:   Tested 1024, 2048, 4096   →  Selected: 2048\n\n"
        "Temperature 0.3: Consistent yet natural.\n"
        "0.0 was too repetitive, 0.7 hallucinated numbers."
    )
    add_text_box(slide, Inches(7), Inches(2.0), Inches(5.5), Inches(3.5),
                 hp_text, font_size=13, color=TEXT_LIGHT)
    add_text_box(slide, Inches(0.8), Inches(4.5), Inches(5), Inches(0.5),
                 "Cost Analysis", font_size=18, bold=True, color=ACCENT_PURPLE)
    add_bullet_points(slide, Inches(0.8), Inches(5.0), Inches(5.5), Inches(2), [
        "~$0.01 per user query",
        "~$0.001 for intent classification alone",
        "Semantic caching could reduce by 40%"
    ])

    # =========== SLIDE 11: LangSmith ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                 "LangSmith Monitoring", font_size=36, bold=True, color=TEXT_WHITE)
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5), Inches(0.5),
                 "What We Track", font_size=18, bold=True, color=ACCENT_PURPLE)
    add_bullet_points(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4), [
        "Every LLM call (intent, RAG, generation)",
        "End-to-end latency per request",
        "Token consumption (input/output)",
        "Error rates and failure modes",
        "Cost per conversation",
        "Full trace of multi-step workflows"
    ])
    add_text_box(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.5),
                 "Insights from Traces", font_size=18, bold=True, color=ACCENT_PURPLE)
    add_bullet_points(slide, Inches(7), Inches(2.0), Inches(5.5), Inches(4), [
        "Intent classification: ~0.8s (fast)",
        "RAG retrieval: ~0.3s (local DB)",
        "Response generation: ~4s (main bottleneck)",
        "Total: ~5.5s average",
        "No failed calls in evaluation run"
    ])

    # =========== SLIDE 12: Trade-offs ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                 "Key Trade-offs", font_size=36, bold=True, color=TEXT_WHITE)
    tradeoffs = (
        "Decision                      Rationale                        Trade-off\n"
        "─────────────────────────────────────────────────────────────────────────\n"
        "LangGraph over CrewAI         Granular routing control         More boilerplate\n"
        "ChromaDB over Pinecone        Zero cost, local                 Not production-scalable\n"
        "GPT-4o over Claude            Better vision + JSON             Higher cost/call\n"
        "Next.js over Streamlit        Professional UI                  Longer development\n"
        "800-char chunks               Context vs precision balance     May split sections"
    )
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(3),
                 tradeoffs, font_size=12, color=TEXT_LIGHT)
    add_text_box(slide, Inches(0.8), Inches(4.5), Inches(11), Inches(0.5),
                 "Hypotheses That Didn't Hold", font_size=18, bold=True, color=ACCENT_PURPLE)
    add_bullet_points(slide, Inches(0.8), Inches(5.0), Inches(11), Inches(2), [
        '"GPT-4o-mini is good enough for everything" → Quality dropped for financial advice',
        '"Users always ask clear questions" → Needed robust intent handling for ambiguous queries'
    ])

    # =========== SLIDE 13: Summary ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                 "Summary", font_size=36, bold=True, color=TEXT_WHITE)
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.7),
                 "FinMate turns transaction data into personalized, actionable financial advice through AI.",
                 font_size=18, color=ACCENT_PURPLE)
    add_text_box(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(0.5),
                 "Technical Highlights", font_size=18, bold=True, color=ACCENT_PURPLE)
    add_bullet_points(slide, Inches(0.8), Inches(2.7), Inches(5.5), Inches(4), [
        "LangGraph workflow — 5 specialist nodes",
        "RAG grounded in curated knowledge",
        "GPT-4o Vision for receipt OCR",
        "Custom MCP server (3 tools)",
        "80% intent accuracy, 3.93/5 relevance",
        "Full LangSmith observability",
        "PII guardrails + injection protection"
    ])
    add_text_box(slide, Inches(7), Inches(2.2), Inches(5.5), Inches(0.5),
                 "Target User & Future", font_size=18, bold=True, color=ACCENT_PURPLE)
    add_bullet_points(slide, Inches(7), Inches(2.7), Inches(5.5), Inches(4), [
        "Young professionals (25-35) who want to manage money",
        "Future: Real bank API (Plaid) integration",
        "Future: Voice interface",
        "Future: Mobile app",
        "Future: Multi-user with auth"
    ])

    # =========== SLIDE 14: Q&A ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text_box(slide, Inches(1), Inches(2.5), Inches(11.3), Inches(1.2),
                 "Q&A", font_size=60, bold=True, color=ACCENT_INDIGO,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(4.0), Inches(11.3), Inches(0.6),
                 "Thank you!", font_size=24, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(5.0), Inches(11.3), Inches(0.5),
                 "github.com/kassymashim/finmate", font_size=14,
                 color=RGBColor(0x64, 0x74, 0x8B), alignment=PP_ALIGN.CENTER)

    # =========== EMPTY SLIDES FOR SCREENSHOTS ===========
    for i in range(4):
        slide = prs.slides.add_slide(blank_layout)
        set_slide_bg(slide)
        add_text_box(slide, Inches(1), Inches(0.5), Inches(11.3), Inches(0.8),
                     f"Screenshot {i + 1}", font_size=28, bold=True, color=TEXT_WHITE,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(1), Inches(3.3), Inches(11.3), Inches(0.8),
                     "[Insert your screenshot here]", font_size=18,
                     color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(1), Inches(4.2), Inches(11.3), Inches(0.5),
                     "Right-click → Insert Picture to add your image", font_size=12,
                     color=RGBColor(0x64, 0x74, 0x8B), alignment=PP_ALIGN.CENTER)

    # Save
    output_path = os.path.join(os.path.dirname(__file__), "FinMate_Presentation.pptx")
    prs.save(output_path)
    print(f"PowerPoint saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)} (14 content + 4 empty for screenshots)")


if __name__ == "__main__":
    create_presentation()
